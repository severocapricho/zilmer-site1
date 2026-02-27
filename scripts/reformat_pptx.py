"""
Script para reformatar PowerPoint: extrai conteúdo e cria nova apresentação com design moderno
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import sys

# Cores da identidade Zilmer
COLORS = {
    'primary': RGBColor(0, 51, 102),      # Azul escuro (#003366)
    'secondary': RGBColor(0, 102, 204),   # Azul médio (#0066cc)
    'accent': RGBColor(0, 51, 102),       # Azul escuro (accent)
    'text_dark': RGBColor(33, 33, 33),    # Cinza escuro
    'text_light': RGBColor(102, 102, 102), # Cinza claro
    'background': RGBColor(255, 255, 255), # Branco
    'background_light': RGBColor(245, 245, 245) # Cinza muito claro
}

def extract_content_from_slide(slide):
    """Extrai todo o conteúdo de um slide"""
    content = {
        'title': '',
        'text_boxes': [],
        'images': [],
        'shapes': []
    }
    
    for shape in slide.shapes:
        # Títulos (geralmente são os primeiros text boxes ou shapes com texto)
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            
            # Identificar títulos (geralmente maiores ou no topo)
            if not content['title'] and len(text) < 100:
                # Verificar se é título (texto curto, pode estar em shape específico)
                if shape.shape_type == 1:  # MSO_SHAPE_TYPE.PLACEHOLDER
                    content['title'] = text
                elif not content['title']:
                    content['title'] = text
            else:
                content['text_boxes'].append({
                    'text': text,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })
        
        # Imagens
        if hasattr(shape, "image"):
            try:
                image_stream = shape.image.stream
                image_bytes = image_stream.read()
                content['images'].append({
                    'bytes': image_bytes,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'ext': shape.image.ext
                })
            except:
                pass
    
    return content

def create_modern_slide(prs, slide_content, slide_index, is_title_slide=False):
    """Cria um slide moderno com o conteúdo extraído"""
    
    # Layout baseado no tipo de slide
    if is_title_slide or slide_index == 0:
        # Slide de título
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS['primary']
        
        # Título principal
        if slide_content['title']:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2), Inches(9), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_content['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(54)
            title_para.font.bold = True
            title_para.font.color.rgb = COLORS['background']
            title_para.alignment = PP_ALIGN.CENTER
            
        # Subtítulo/descrição
        if slide_content['text_boxes']:
            subtitle_text = slide_content['text_boxes'][0]['text']
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4), Inches(8), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.name = 'Calibri'
    
    else:
        # Slide de conteúdo
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLORS['background']
        
        # Barra lateral colorida
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = COLORS['primary']
        left_bar.line.fill.background()
        
        # Título do slide
        if slide_content['title']:
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.5), Inches(8.5), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_content['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = COLORS['primary']
            title_para.font.name = 'Calibri'
        
        # Conteúdo de texto
        content_y = Inches(1.8)
        for i, text_box in enumerate(slide_content['text_boxes']):
            if text_box['text'] and text_box['text'] != slide_content['title']:
                content_box = slide.shapes.add_textbox(
                    Inches(0.8), content_y, Inches(8.5), Inches(1)
                )
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                content_frame.text = text_box['text']
                
                for para in content_frame.paragraphs:
                    para.font.size = Pt(18)
                    para.font.color.rgb = COLORS['text_dark']
                    para.font.name = 'Calibri'
                    para.space_after = Pt(12)
                
                content_y += Inches(1.2)
        
        # Imagens (lado direito ou abaixo)
        if slide_content['images']:
            img_y = Inches(1.8)
            img_x = Inches(5.5)
            for img_data in slide_content['images']:
                try:
                    slide.shapes.add_picture(
                        img_data['bytes'], img_x, img_y, 
                        width=Inches(4), height=Inches(4.5)
                    )
                except:
                    pass

def reformat_presentation(input_path, output_path):
    """Função principal para reformatar a apresentação"""
    
    if not os.path.exists(input_path):
        print(f"❌ Arquivo não encontrado: {input_path}")
        return False
    
    print(f"📖 Lendo apresentação: {input_path}")
    
    try:
        # Ler apresentação original
        prs_original = Presentation(input_path)
        
        # Criar nova apresentação
        prs_new = Presentation()
        prs_new.slide_width = Inches(10)
        prs_new.slide_height = Inches(7.5)
        
        print(f"📊 Total de slides: {len(prs_original.slides)}")
        
        # Processar cada slide
        for i, slide in enumerate(prs_original.slides):
            print(f"  Processando slide {i+1}/{len(prs_original.slides)}...")
            
            # Extrair conteúdo
            content = extract_content_from_slide(slide)
            
            # Criar novo slide com design moderno
            is_title = (i == 0)
            create_modern_slide(prs_new, content, i, is_title_slide=is_title)
        
        # Salvar nova apresentação
        prs_new.save(output_path)
        print(f"✅ Nova apresentação salva: {output_path}")
        return True
        
    except Exception as e:
        print(f"❌ Erro: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python reformat_pptx.py <caminho_do_pptx> [caminho_saida]")
        print("\nExemplo:")
        print("  python reformat_pptx.py apresentacao.pptx")
        print("  python reformat_pptx.py apresentacao.pptx nova_apresentacao.pptx")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else input_file.replace('.pptx', '_reformatado.pptx')
    
    reformat_presentation(input_file, output_file)

